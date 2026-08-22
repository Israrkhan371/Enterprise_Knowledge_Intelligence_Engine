"""
Integration tests: for every knowledge source type, run the REAL loader on
a small generated fixture, chunk it with the REAL chunker, store it via a
REAL Postgres write, and confirm keyword_search() finds a unique marker
word planted in that source's content.

This is deliberately narrower than the full ingestion pipeline
(app.ingestion.pipeline.ingest_document): it skips embeddings/ChromaDB and
Neo4j graph population, which need API credentials / a running Neo4j
instance and aren't what this suite is checking. What's under test here is
specifically "does keyword search work against everything the loaders
produce" — the same question the manual terminal session answered once,
by hand, non-repeatably. This makes that check automated and run on every
test invocation instead.

Requires a real Postgres connection (the same one the app itself uses via
DATABASE_URL / app.core.config.settings). If it can't connect, the whole
module is skipped rather than erroring, so this doesn't break environments
without a DB (e.g. a laptop running only tests/test_loaders.py).

Each test runs inside a transaction that is ALWAYS rolled back in the `db`
fixture's teardown — nothing here is ever committed, so no test data is
left behind in a real dev database. keyword_search() still sees the
uncommitted rows because it runs its raw SQL on the same session/connection
within the same transaction.
"""
import json
import uuid

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from sqlalchemy import text as sql_text

from app.core.database import SessionLocal
from app.core.models import Document, DocumentChunk
from app.ingestion.chunking import chunk_text
from app.ingestion.loaders import load_by_source_type
from app.search.keyword import keyword_search


@pytest.fixture
def db():
    session = SessionLocal()
    try:
        session.execute(sql_text("SELECT 1"))
    except Exception as exc:
        session.close()
        pytest.skip(f"No live Postgres connection available for integration tests: {exc}")
    yield session
    session.rollback()  # discard everything written during the test
    session.close()


def _marker() -> str:
    """A unique, unambiguous token per test run/call, so tests can't pass
    by matching leftover data from a previous (failed) run or another
    test's fixture."""
    return f"marker{uuid.uuid4().hex[:12]}"


def _store_and_search(db, source_type: str, text: str, marker: str) -> list[dict]:
    """Shared tail: real chunker -> real DocumentChunk rows -> real keyword_search."""
    document = Document(title=f"test-{source_type}", source_type=source_type)
    db.add(document)
    db.flush()  # assigns document.id without committing, see pipeline.py's note on why

    chunks = chunk_text(text)
    assert chunks, f"chunk_text() produced no chunks for source_type={source_type!r}"

    for idx, chunk in enumerate(chunks):
        db.add(DocumentChunk(document_id=document.id, chunk_index=idx, text=chunk))
    db.flush()

    return keyword_search(db, marker, top_k=5)


# --- one test per source type -------------------------------------------

def test_keyword_search_finds_markdown_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.md"
    path.write_text(f"# Notes\n\nProject uses {marker} as a codeword.\n", encoding="utf-8")

    text = load_by_source_type("markdown", str(path))
    results = _store_and_search(db, "markdown", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_code_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.py"
    path.write_text(f"# module\ndef f():\n    return '{marker}'\n", encoding="utf-8")

    text = load_by_source_type("code", str(path))
    results = _store_and_search(db, "code", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_transcript_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.vtt"
    path.write_text(
        f"WEBVTT\n\n00:00:01.000 --> 00:00:04.000\nWe discussed {marker} today.\n",
        encoding="utf-8",
    )

    text = load_by_source_type("transcript", str(path))
    results = _store_and_search(db, "transcript", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_meeting_notes_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.txt"
    path.write_text(f"Meeting notes: {marker} is the action item.\n", encoding="utf-8")

    text = load_by_source_type("meeting_notes", str(path))
    results = _store_and_search(db, "meeting_notes", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_blog_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.html"
    path.write_text(f"<html><body><p>This post is about {marker}.</p></body></html>", encoding="utf-8")

    text = load_by_source_type("blog", str(path))
    results = _store_and_search(db, "blog", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_api_docs_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.json"
    spec = {"paths": {f"/{marker}": {"get": {"summary": "Test endpoint"}}}}
    path.write_text(json.dumps(spec), encoding="utf-8")

    text = load_by_source_type("api_docs", str(path))
    # api_docs content includes the marker with a leading slash (it's a URL
    # path segment) — see the flangeburst987 tokenization note in
    # app/search/keyword.py; search WITH the slash, matching how the path
    # actually appears in the endpoint text, same as the manual session did.
    results = _store_and_search(db, "api_docs", text, f"/{marker}")

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_db_schema_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.sql"
    path.write_text(f"CREATE TABLE {marker} (id INT PRIMARY KEY);\n", encoding="utf-8")

    text = load_by_source_type("db_schema", str(path))
    results = _store_and_search(db, "db_schema", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_lms_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.html"
    path.write_text(f"<html><body><h1>Lesson</h1><p>Topic: {marker}</p></body></html>", encoding="utf-8")

    text = load_by_source_type("lms", str(path))
    results = _store_and_search(db, "lms", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_docx_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.docx"
    doc = DocxDocument()
    doc.add_paragraph(f"This document mentions {marker} in the body text.")
    doc.save(str(path))

    text = load_by_source_type("docx", str(path))
    results = _store_and_search(db, "docx", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_xlsx_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = f"This spreadsheet mentions {marker} in a cell."
    wb.save(str(path))

    text = load_by_source_type("xlsx", str(path))
    results = _store_and_search(db, "xlsx", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_pptx_content(db, tmp_path):
    marker = _marker()
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = f"This slide mentions {marker} in the body."
    prs.save(str(path))

    text = load_by_source_type("pptx", str(path))
    results = _store_and_search(db, "pptx", text, marker)

    assert any(marker in r["text"] for r in results)


def test_keyword_search_finds_github_content(db, monkeypatch):
    """
    load_github_repo() makes a real network call, so it's mocked here (same
    approach as tests/test_loaders.py's GitHub tests) rather than depending
    on network access or a token — this closes the "GitHub never confirmed
    with keyword search" gap without depending on the actual 403 rate-limit
    issue seen in the live terminal session being resolved first.
    """
    from app.ingestion import pipeline

    marker = _marker()
    fake_files = [{"path": "README.md", "content": f"This repo is about {marker}.", "sha": "abc123"}]
    monkeypatch.setattr(pipeline, "load_github_repo", lambda repo_url, github_token=None: fake_files)

    files = pipeline.load_github_repo("https://github.com/example/example")
    text = files[0]["content"]

    results = _store_and_search(db, "github", text, marker)

    assert any(marker in r["text"] for r in results)


@pytest.mark.skip(
    reason=(
        "Real-text PDF extraction depends on unstructured's hi_res layout "
        "model / OCR stack, which test_loaders.py's own sample_pdf_path "
        "fixture deliberately doesn't exercise (see its docstring) — a "
        "hand-built PDF content stream isn't a reliable stand-in. Requires "
        "a real sample.pdf with known text; wire this up once "
        "tests/fixtures/sample.pdf exists, same as "
        "test_load_pdf_on_real_sample_if_present in test_loaders.py."
    )
)
def test_keyword_search_finds_pdf_content(db):
    pass


def test_keyword_search_returns_empty_for_term_not_in_any_source(db):
    """Sanity check that the marker-based tests above aren't passing by
    accident (e.g. keyword_search returning everything regardless of
    query) — a term that was never planted anywhere should find nothing."""
    marker = _marker()
    document = Document(title="unrelated", source_type="markdown")
    db.add(document)
    db.flush()
    db.add(DocumentChunk(document_id=document.id, chunk_index=0, text="completely unrelated content"))
    db.flush()

    results = keyword_search(db, marker, top_k=5)

    assert results == []
