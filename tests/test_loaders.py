"""
Tests for app/ingestion/loaders.py.

Deliberately self-contained: instead of depending on personal sample files
outside the repo (which won't exist on a fresh clone / CI runner / Engineer
B's machine), each PDF/docx fixture is generated on the fly with pypdf /
python-docx — both already in requirements-lock.txt as unstructured deps.
"""
import pytest
from pypdf import PdfWriter
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from app.ingestion.loaders import (
    load_pdf,
    load_docx,
    load_xlsx,
    load_pptx,
    load_markdown,
    load_code,
    load_transcript,
    load_meeting_notes,
    load_blog,
    load_github_repo,
    load_api_docs,
    load_db_schema,
    load_lms,
    load_by_source_type,
    SOURCE_LOADERS,
)

EXPECTED_TEXT = "EKIE loader test fixture — hello from pytest."


@pytest.fixture
def sample_pdf_path(tmp_path):
    path = tmp_path / "sample.pdf"
    # Build a minimal valid PDF with real text via pypdf's low-level API
    # (avoids adding reportlab as a new dependency just for this fixture).
    writer = PdfWriter()
    page = writer.add_blank_page(width=300, height=150)
    # pypdf's blank page has no text layer; unstructured's hi_res strategy
    # still runs its layout model over the (blank) page image and returns
    # an empty element list rather than erroring, so we assert on "ran
    # without raising and returned a str" here, and cover real extracted
    # text in test_load_pdf_on_real_sample_if_present below instead.
    writer.write(str(path))
    return str(path)


@pytest.fixture
def sample_docx_path(tmp_path):
    path = tmp_path / "sample.docx"
    doc = DocxDocument()
    doc.add_paragraph(EXPECTED_TEXT)
    doc.save(str(path))
    return str(path)


@pytest.fixture
def sample_xlsx_path(tmp_path):
    path = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = EXPECTED_TEXT
    wb.save(str(path))
    return str(path)


@pytest.fixture
def sample_pptx_path(tmp_path):
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content layout
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = EXPECTED_TEXT
    prs.save(str(path))
    return str(path)


@pytest.fixture
def sample_markdown_path(tmp_path):
    path = tmp_path / "sample.md"
    path.write_text(f"# Heading\n\n{EXPECTED_TEXT}\n", encoding="utf-8")
    return str(path)


@pytest.fixture
def sample_code_path(tmp_path):
    path = tmp_path / "sample.py"
    path.write_text(f"# {EXPECTED_TEXT}\ndef f():\n    return 1\n", encoding="utf-8")
    return str(path)


# --- docx -------------------------------------------------------------

def test_load_docx_returns_expected_text(sample_docx_path):
    text = load_docx(sample_docx_path)
    assert EXPECTED_TEXT in text


def test_load_docx_raises_on_missing_file():
    with pytest.raises(Exception):
        load_docx("/nonexistent/path/sample.docx")


# --- xlsx ("Office Files" per the case-study brief) ---------------------

def test_load_xlsx_returns_expected_text(sample_xlsx_path):
    text = load_xlsx(sample_xlsx_path)
    assert EXPECTED_TEXT in text


def test_load_xlsx_raises_on_missing_file():
    with pytest.raises(Exception):
        load_xlsx("/nonexistent/path/sample.xlsx")


def test_load_xlsx_raises_on_empty_workbook(tmp_path):
    # An all-empty sheet is a real case (a template someone uploads by
    # mistake, or a sheet that's genuinely just formatting/no data) -
    # should fail loudly via ValueError, not silently ingest a
    # zero-content document that then can't ever surface in search.
    path = tmp_path / "empty.xlsx"
    Workbook().save(str(path))
    with pytest.raises(ValueError, match="No extractable content"):
        load_xlsx(str(path))


def test_load_xlsx_is_registered_in_source_loaders():
    assert SOURCE_LOADERS["xlsx"] is load_xlsx


# --- pptx ("Office Files" per the case-study brief) ----------------------

def test_load_pptx_returns_expected_text(sample_pptx_path):
    text = load_pptx(sample_pptx_path)
    assert EXPECTED_TEXT in text


def test_load_pptx_keeps_slide_titles(sample_pptx_path):
    # Unlike load_docx (which filters Title elements out - see its
    # docstring), load_pptx deliberately keeps them: a slide's bullets
    # are often meaningless without the title they're under.
    text = load_pptx(sample_pptx_path)
    assert "Test Slide" in text


def test_load_pptx_raises_on_missing_file():
    with pytest.raises(Exception):
        load_pptx("/nonexistent/path/sample.pptx")


def test_load_pptx_raises_on_empty_presentation(tmp_path):
    path = tmp_path / "empty.pptx"
    Presentation().save(str(path))  # zero slides
    with pytest.raises(ValueError, match="No extractable content"):
        load_pptx(str(path))


def test_load_pptx_is_registered_in_source_loaders():
    assert SOURCE_LOADERS["pptx"] is load_pptx


# --- pdf ----------------------------------------------------------------
# Generating a PDF with a real extractable text layer needs a PDF-writing
# library beyond pypdf's blank-page helper, so this suite focuses on "runs
# without raising and returns a str" for the synthetic fixture, plus an
# opt-in check against a real sample if one is dropped in tests/fixtures/
# (matches how you Level-1 tested this manually with sample.pdf).

def test_load_pdf_runs_without_raising(sample_pdf_path):
    text = load_pdf(sample_pdf_path)
    assert isinstance(text, str)


def test_load_pdf_raises_on_missing_file():
    with pytest.raises(Exception):
        load_pdf("/nonexistent/path/sample.pdf")


def test_load_pdf_on_real_sample_if_present():
    """
    Optional stronger check: drop a real PDF at tests/fixtures/sample.pdf
    (e.g. copy the one you Level-1 tested manually) and this asserts real
    text actually comes out, not just "didn't crash".
    """
    from pathlib import Path
    fixture = Path(__file__).parent / "fixtures" / "sample.pdf"
    if not fixture.exists():
        pytest.skip("tests/fixtures/sample.pdf not present — skipping real-text check")
    text = load_pdf(str(fixture))
    assert len(text) > 0


# --- markdown -------------------------------------------------------------

def test_load_markdown_returns_expected_text(sample_markdown_path):
    text = load_markdown(sample_markdown_path)
    assert EXPECTED_TEXT in text
    assert "# Heading" in text


def test_load_markdown_raises_on_missing_file():
    # Path.read_text on a missing file raises FileNotFoundError, not a
    # silent empty string — confirms load_markdown doesn't swallow errors.
    with pytest.raises(FileNotFoundError):
        load_markdown("/nonexistent/path/sample.md")


# --- code -------------------------------------------------------------

def test_load_code_returns_expected_text(sample_code_path):
    text = load_code(sample_code_path)
    assert EXPECTED_TEXT in text
    assert "def f():" in text


def test_load_code_includes_filename_header(sample_code_path):
    text = load_code(sample_code_path)
    assert text.startswith("# File: sample.py")


# --- transcript -------------------------------------------------------

def test_load_transcript_plain_text_passthrough(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text(EXPECTED_TEXT, encoding="utf-8")
    text = load_transcript(str(path))
    assert EXPECTED_TEXT in text


def test_load_transcript_strips_vtt_timestamps_and_header(tmp_path):
    # A realistic WebVTT export: header, cue timestamps, spoken lines.
    vtt_content = (
        "WEBVTT\n\n"
        "00:00:01.000 --> 00:00:04.000\n"
        f"{EXPECTED_TEXT}\n\n"
        "00:00:04.500 --> 00:00:07.200\n"
        "Second line of dialogue.\n"
    )
    path = tmp_path / "sample.vtt"
    path.write_text(vtt_content, encoding="utf-8")

    text = load_transcript(str(path))

    assert EXPECTED_TEXT in text
    assert "Second line of dialogue." in text
    assert "WEBVTT" not in text
    assert "-->" not in text
    assert "00:00:01.000" not in text


def test_load_transcript_strips_srt_cue_numbers_and_timestamps(tmp_path):
    # A realistic SRT export: numeric cue index, comma-decimal timestamps.
    srt_content = (
        "1\n"
        "00:00:01,000 --> 00:00:04,000\n"
        f"{EXPECTED_TEXT}\n\n"
        "2\n"
        "00:00:04,500 --> 00:00:07,200\n"
        "Second line of dialogue.\n"
    )
    path = tmp_path / "sample.srt"
    path.write_text(srt_content, encoding="utf-8")

    text = load_transcript(str(path))

    assert EXPECTED_TEXT in text
    assert "Second line of dialogue." in text
    assert "-->" not in text
    # bare cue-number lines ("1", "2") should be stripped, not just timestamps
    lines = text.splitlines()
    assert "1" not in lines
    assert "2" not in lines


def test_load_transcript_raises_on_missing_file():
    with pytest.raises(FileNotFoundError):
        load_transcript("/nonexistent/path/sample.vtt")


# --- meeting_notes -------------------------------------------------------

def test_load_meeting_notes_plain_text(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text(f"Action items:\n{EXPECTED_TEXT}\n", encoding="utf-8")
    text = load_meeting_notes(str(path))
    assert EXPECTED_TEXT in text
    assert "Action items:" in text


def test_load_meeting_notes_strips_markup_if_present(tmp_path):
    # Notes exported from a bot alongside timestamps should still get
    # cleaned, same as a real transcript.
    content = (
        "00:00:01.000 --> 00:00:04.000\n"
        f"{EXPECTED_TEXT}\n"
    )
    path = tmp_path / "notes_with_timestamps.txt"
    path.write_text(content, encoding="utf-8")
    text = load_meeting_notes(str(path))
    assert EXPECTED_TEXT in text
    assert "-->" not in text


def test_load_meeting_notes_raises_on_missing_file():
    with pytest.raises(FileNotFoundError):
        load_meeting_notes("/nonexistent/path/notes.txt")


# --- blog ---------------------------------------------------------------

def test_load_blog_from_local_html_file(tmp_path):
    html = (
        "<html><head><style>body{color:red}</style>"
        "<script>console.log('noise')</script></head>"
        f"<body><h1>Title</h1><p>{EXPECTED_TEXT}</p></body></html>"
    )
    path = tmp_path / "post.html"
    path.write_text(html, encoding="utf-8")

    text = load_blog(str(path))

    assert EXPECTED_TEXT in text
    assert "Title" in text
    # script/style content must not leak into extracted text
    assert "console.log" not in text
    assert "color:red" not in text


def test_load_blog_from_url_mocked(monkeypatch):
    import httpx

    class FakeResponse:
        def __init__(self, text):
            self.text = text

        def raise_for_status(self):
            pass

    class FakeClient:
        def __init__(self, *args, **kwargs):
            pass

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

        def get(self, url):
            html = f"<html><body><p>{EXPECTED_TEXT}</p></body></html>"
            return FakeResponse(html)

    monkeypatch.setattr(httpx, "Client", FakeClient)

    text = load_blog("https://example.com/blog/post")
    assert EXPECTED_TEXT in text


def test_load_blog_raises_on_missing_local_file():
    with pytest.raises(FileNotFoundError):
        load_blog("/nonexistent/path/post.html")


# --- github_repo (mocked — no real network call) -----------------------

def test_load_github_repo_returns_list_of_dicts(monkeypatch):
    """
    load_github_repo hits the real GitHub API, so this mocks httpx.Client
    entirely rather than making a live network call in tests. This mock
    simulates a repo with a nested folder structure to confirm recursion
    works, not just flat root-level files — a real bug (this loader only
    ever saw README.md at the repo root) was caught by manual testing
    against a real multi-file repo and fixed here; this mock reproduces
    that same shape so it can't silently regress again.
    """
    import httpx

    class FakeResponse:
        def __init__(self, json_data=None, text=""):
            self._json = json_data
            self.text = text

        def raise_for_status(self):
            pass

        def json(self):
            return self._json

    # Simulated tree:
    #   README.md
    #   image.png                (wrong extension -> excluded)
    #   app/                     (subfolder -> must recurse into it)
    #     main.py
    #   node_modules/            (excluded directory -> must NOT recurse)
    #     junk.md
    responses = {
        "https://api.github.com/repos/example/repo/contents": [
            {"type": "file", "name": "README.md", "path": "README.md",
             "download_url": "https://example.com/README.md"},
            {"type": "file", "name": "image.png", "path": "image.png",
             "download_url": "https://example.com/image.png"},
            {"type": "dir", "name": "app",
             "url": "https://api.github.com/repos/example/repo/contents/app"},
            {"type": "dir", "name": "node_modules",
             "url": "https://api.github.com/repos/example/repo/contents/node_modules"},
        ],
        "https://api.github.com/repos/example/repo/contents/app": [
            {"type": "file", "name": "main.py", "path": "app/main.py",
             "download_url": "https://example.com/main.py"},
        ],
    }

    class FakeClient:
        def __init__(self, *args, **kwargs):
            pass

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

        def get(self, url):
            if url in responses:
                return FakeResponse(json_data=responses[url])
            if url == "https://api.github.com/repos/example/repo/contents/node_modules":
                raise AssertionError(
                    "load_github_repo must not descend into excluded directories"
                )
            return FakeResponse(text=EXPECTED_TEXT)

    monkeypatch.setattr(httpx, "Client", FakeClient)

    results = load_github_repo("https://github.com/example/repo")

    paths = {r["path"] for r in results}
    assert paths == {"README.md", "app/main.py"}
    # image.png (wrong extension) and node_modules/junk.md (excluded dir)
    # must both be absent
    assert "image.png" not in paths
    assert not any("node_modules" in p for p in paths)


def test_load_github_repo_respects_max_depth(monkeypatch):
    """
    Confirms recursion actually stops past _GITHUB_MAX_DEPTH rather than
    walking an arbitrarily deep tree forever.
    """
    import httpx
    from app.ingestion import loaders as loaders_module

    call_count = {"n": 0}

    class FakeResponse:
        def __init__(self, json_data):
            self._json = json_data

        def raise_for_status(self):
            pass

        def json(self):
            return self._json

    class FakeClient:
        def __init__(self, *args, **kwargs):
            pass

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

        def get(self, url):
            call_count["n"] += 1
            # Every folder contains exactly one subfolder, forever —
            # without a depth guard this would recurse indefinitely.
            return FakeResponse(json_data=[
                {"type": "dir", "name": "nested", "url": url + "/nested"},
            ])

    monkeypatch.setattr(httpx, "Client", FakeClient)

    results = load_github_repo("https://github.com/example/infinite-repo")

    assert results == []
    assert call_count["n"] <= loaders_module._GITHUB_MAX_DEPTH + 1


# --- api_docs (OpenAPI/Swagger JSON) ------------------------------------

def test_load_api_docs_includes_title_and_endpoint(tmp_path):
    import json
    spec = {
        "info": {"title": "EKIE API", "version": "1.0.0", "description": EXPECTED_TEXT},
        "paths": {
            "/documents": {
                "get": {
                    "summary": "List documents",
                    "parameters": [{"name": "category_id", "in": "query", "required": False}],
                    "responses": {"200": {"description": "A list of documents"}},
                }
            }
        },
    }
    path = tmp_path / "openapi.json"
    path.write_text(json.dumps(spec), encoding="utf-8")

    text = load_api_docs(str(path))

    assert "EKIE API" in text
    assert EXPECTED_TEXT in text
    assert "GET /documents" in text
    assert "category_id" in text
    assert "200: A list of documents" in text


def test_load_api_docs_skips_non_http_keys(tmp_path):
    # "parameters" can appear as a path-level sibling key (shared across
    # methods) rather than only inside a method — must not be treated as
    # if it were an HTTP method itself.
    import json
    spec = {
        "info": {"title": "EKIE API"},
        "paths": {
            "/documents": {
                "parameters": [{"name": "shared_param", "in": "query"}],
                "get": {"summary": "List documents"},
            }
        },
    }
    path = tmp_path / "openapi.json"
    path.write_text(json.dumps(spec), encoding="utf-8")

    text = load_api_docs(str(path))
    assert "GET /documents" in text


def test_load_api_docs_raises_on_missing_file():
    with pytest.raises(FileNotFoundError):
        load_api_docs("/nonexistent/path/openapi.json")


# --- db_schema (.sql dumps) ----------------------------------------------

def test_load_db_schema_keeps_ddl(tmp_path):
    sql = (
        "CREATE TABLE documents (\n"
        "    id UUID PRIMARY KEY,\n"
        f"    title TEXT -- {EXPECTED_TEXT}\n"
        ");\n"
        "ALTER TABLE documents ADD CONSTRAINT fk_category "
        "FOREIGN KEY (category_id) REFERENCES categories(id);\n"
    )
    path = tmp_path / "schema.sql"
    path.write_text(sql, encoding="utf-8")

    text = load_db_schema(str(path))

    assert "CREATE TABLE documents" in text
    assert "ALTER TABLE documents" in text
    assert EXPECTED_TEXT in text


def test_load_db_schema_drops_insert_statements(tmp_path):
    sql = (
        "CREATE TABLE documents (id UUID PRIMARY KEY);\n"
        "INSERT INTO documents (id) VALUES ('11111111-1111-1111-1111-111111111111');\n"
    )
    path = tmp_path / "schema.sql"
    path.write_text(sql, encoding="utf-8")

    text = load_db_schema(str(path))

    assert "CREATE TABLE documents" in text
    assert "INSERT INTO" not in text
    assert "11111111" not in text


def test_load_db_schema_drops_copy_data_block(tmp_path):
    sql = (
        "CREATE TABLE documents (id UUID PRIMARY KEY, title TEXT);\n"
        "COPY documents (id, title) FROM stdin;\n"
        "11111111-1111-1111-1111-111111111111\tSecret internal title\n"
        "\\.\n"
        "CREATE INDEX idx_documents_title ON documents (title);\n"
    )
    path = tmp_path / "schema.sql"
    path.write_text(sql, encoding="utf-8")

    text = load_db_schema(str(path))

    assert "CREATE TABLE documents" in text
    assert "CREATE INDEX idx_documents_title" in text
    assert "Secret internal title" not in text
    assert "COPY documents" not in text


def test_load_db_schema_raises_on_missing_file():
    with pytest.raises(FileNotFoundError):
        load_db_schema("/nonexistent/path/schema.sql")


# --- lms (SCORM zip / exported HTML) --------------------------------------

def test_load_lms_from_html_file(tmp_path):
    html = f"<html><body><h1>Lesson 1</h1><p>{EXPECTED_TEXT}</p></body></html>"
    path = tmp_path / "lesson.html"
    path.write_text(html, encoding="utf-8")

    text = load_lms(str(path))

    assert "Lesson 1" in text
    assert EXPECTED_TEXT in text


def test_load_lms_from_scorm_zip(tmp_path):
    import zipfile

    zip_path = tmp_path / "course.zip"
    with zipfile.ZipFile(zip_path, "w") as archive:
        archive.writestr(
            "imsmanifest.xml",
            "<manifest><organizations/></manifest>",
        )
        archive.writestr(
            "lesson1.html",
            f"<html><body><h1>Lesson 1</h1><p>{EXPECTED_TEXT}</p></body></html>",
        )
        archive.writestr(
            "lesson2.html",
            "<html><body><h1>Lesson 2</h1><p>Second lesson content.</p></body></html>",
        )

    text = load_lms(str(zip_path))

    assert "lesson1.html" in text
    assert EXPECTED_TEXT in text
    assert "lesson2.html" in text
    assert "Second lesson content." in text
    # manifest is packaging metadata, not course content — must not appear
    assert "organizations" not in text


def test_load_lms_raises_on_missing_file():
    with pytest.raises(FileNotFoundError):
        load_lms("/nonexistent/path/lesson.html")


# --- dispatch table / load_by_source_type ------------------------------

def test_source_loaders_registry_has_expected_keys():
    assert set(SOURCE_LOADERS.keys()) == {
        "pdf", "docx", "xlsx", "pptx", "markdown", "code", "transcript",
        "meeting_notes", "blog", "api_docs", "db_schema", "lms",
    }


def test_github_excluded_from_source_loaders():
    # Documented on purpose: load_github_repo returns list[dict], not str,
    # so it can't go through load_by_source_type()/ingest_document().
    assert "github" not in SOURCE_LOADERS


def test_load_by_source_type_github_raises_helpful_error():
    with pytest.raises(ValueError, match="ingest_github_repo"):
        load_by_source_type("github", "irrelevant-path")


def test_load_by_source_type_unknown_raises():
    with pytest.raises(ValueError, match="No loader registered"):
        load_by_source_type("carrier-pigeon", "irrelevant-path")


def test_load_by_source_type_dispatches_correctly(sample_markdown_path):
    text = load_by_source_type("markdown", sample_markdown_path)
    assert EXPECTED_TEXT in text
